"""
Orquestra - Transcription & Vision Service
Audio transcription via OpenRouter Whisper (large files) or Groq Whisper (small files),
with ffmpeg chunking as extra fallback. Image description via OpenRouter.
"""

import base64
import logging
import os
import subprocess
import tempfile

import httpx
import fitz  # pymupdf

from app.config import settings

logger = logging.getLogger(__name__)

# Groq Whisper limit is 25MB. Use 20min chunks to stay safe.
CHUNK_DURATION_SECONDS = 20 * 60  # 20 minutes
MAX_FILE_SIZE_BYTES = 24 * 1024 * 1024  # 24MB (safe margin under 25MB)


async def _transcribe_single(file_path: str) -> str:
    """Transcribe a single audio file via Groq Whisper API."""
    url = "https://api.groq.com/openai/v1/audio/transcriptions"
    headers = {"Authorization": f"Bearer {settings.GROQ_API_KEY}"}

    filename = os.path.basename(file_path)

    async with httpx.AsyncClient(timeout=180.0) as client:
        with open(file_path, "rb") as audio_file:
            files = {
                "file": (filename, audio_file, "audio/ogg"),
                "model": (None, "whisper-large-v3"),
                "language": (None, "pt"),
            }
            response = await client.post(url, headers=headers, files=files)
            response.raise_for_status()

    data = response.json()
    text = data.get("text", "").strip()
    logger.info("[TRANSCRIBER] Transcribed %s -> %d chars", filename, len(text))
    return text


async def _transcribe_via_openrouter(file_path: str) -> str:
    """Transcribe audio via OpenRouter chat completions with input_audio (no 25MB limit)."""
    url = f"{settings.OPENROUTER_BASE_URL}/chat/completions"
    headers = {
        "Authorization": f"Bearer {settings.OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
    }

    # Read and base64-encode the audio file
    with open(file_path, "rb") as f:
        audio_b64 = base64.b64encode(f.read()).decode("utf-8")

    # Detect format from extension
    ext = os.path.splitext(file_path)[1].lstrip(".").lower()
    fmt_map = {"webm": "ogg", "opus": "ogg", "m4a": "m4a", "mp3": "mp3", "wav": "wav"}
    audio_format = fmt_map.get(ext, "ogg")

    payload = {
        "model": settings.MODEL_TRANSCRIPTION,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": "Transcreva este audio completamente em portugues. Retorne APENAS a transcricao, sem comentarios.",
                    },
                    {
                        "type": "input_audio",
                        "input_audio": {
                            "data": audio_b64,
                            "format": audio_format,
                        },
                    },
                ],
            }
        ],
        "max_tokens": 16000,
        "temperature": 0.0,
    }

    async with httpx.AsyncClient(timeout=300.0) as client:
        response = await client.post(url, headers=headers, json=payload)
        response.raise_for_status()

    data = response.json()
    text = data["choices"][0]["message"]["content"].strip()
    logger.info("[TRANSCRIBER] OpenRouter transcribed %s -> %d chars", os.path.basename(file_path), len(text))
    return text


def _split_audio_chunks(file_path: str, chunk_seconds: int = CHUNK_DURATION_SECONDS) -> list[str]:
    """
    Split audio file into chunks using ffmpeg.
    Returns list of chunk file paths.
    """
    tmp_dir = tempfile.mkdtemp(prefix="orquestra_chunks_")

    # Get audio duration using ffprobe
    try:
        result = subprocess.run(
            [
                "ffprobe", "-v", "error",
                "-show_entries", "format=duration",
                "-of", "default=noprint_wrappers=1:nokey=1",
                file_path,
            ],
            capture_output=True, text=True, timeout=30,
        )
        duration = float(result.stdout.strip())
    except Exception as exc:
        logger.warning("[TRANSCRIBER] ffprobe failed, processing as single file: %s", exc)
        return [file_path]

    if duration <= chunk_seconds:
        return [file_path]

    # Split into chunks
    chunk_paths = []
    chunk_index = 0
    start = 0

    while start < duration:
        chunk_path = os.path.join(tmp_dir, f"chunk_{chunk_index:03d}.ogg")
        try:
            subprocess.run(
                [
                    "ffmpeg", "-y", "-i", file_path,
                    "-ss", str(start),
                    "-t", str(chunk_seconds),
                    "-vn",  # no video
                    "-acodec", "libopus",
                    "-b:a", "48k",  # compress to fit under limit
                    chunk_path,
                ],
                capture_output=True, timeout=120,
            )
            if os.path.exists(chunk_path) and os.path.getsize(chunk_path) > 0:
                chunk_paths.append(chunk_path)
        except Exception as exc:
            logger.error("[TRANSCRIBER] ffmpeg chunk %d failed: %s", chunk_index, exc)

        start += chunk_seconds
        chunk_index += 1

    logger.info(
        "[TRANSCRIBER] Split %.0fs audio into %d chunks",
        duration, len(chunk_paths),
    )
    return chunk_paths if chunk_paths else [file_path]


async def transcribe_audio(file_path: str) -> str:
    """
    Transcribe an audio file.

    Strategy:
    - Small files (≤24MB): Groq Whisper directly (fast).
    - Large files (>24MB): OpenRouter Whisper first (no size limit);
      falls back to ffmpeg chunking via Groq if OpenRouter fails.

    Args:
        file_path: Path to the audio file on disk.

    Returns:
        Transcription text.
    """
    has_groq = bool(settings.GROQ_API_KEY)
    has_openrouter = bool(settings.OPENROUTER_API_KEY)

    if not has_groq and not has_openrouter:
        raise RuntimeError(
            "No transcription API key configured. Set GROQ_API_KEY or OPENROUTER_API_KEY."
        )

    file_size = os.path.getsize(file_path)

    # Small files: Groq Whisper directly (no need to involve OpenRouter)
    if file_size <= MAX_FILE_SIZE_BYTES:
        if has_groq:
            return await _transcribe_single(file_path)
        return await _transcribe_via_openrouter(file_path)

    # Large files: try OpenRouter first (handles files beyond Groq's 25MB limit)
    logger.info(
        "[TRANSCRIBER] File %s is %.1fMB, trying OpenRouter...",
        os.path.basename(file_path),
        file_size / (1024 * 1024),
    )

    if has_openrouter:
        try:
            return await _transcribe_via_openrouter(file_path)
        except Exception as exc:
            logger.warning(
                "[TRANSCRIBER] OpenRouter failed for %.1fMB file: %s. Falling back to chunking...",
                file_size / (1024 * 1024),
                exc,
            )

    # Fallback: split into chunks and transcribe each via Groq
    if not has_groq:
        raise RuntimeError(
            "OpenRouter transcription failed and GROQ_API_KEY is not configured. "
            "Cannot transcribe large file."
        )

    logger.info(
        "[TRANSCRIBER] Splitting %s into chunks...",
        os.path.basename(file_path),
    )
    chunks = _split_audio_chunks(file_path)
    transcriptions = []
    failed_chunks = 0
    tmp_dir = None

    for i, chunk_path in enumerate(chunks):
        try:
            logger.info("[TRANSCRIBER] Transcribing chunk %d/%d...", i + 1, len(chunks))
            text = await _transcribe_single(chunk_path)
            if text:
                transcriptions.append(text)
        except Exception as exc:
            logger.error("[TRANSCRIBER] Chunk %d/%d failed: %s", i + 1, len(chunks), exc)
            failed_chunks += 1
        finally:
            if chunk_path != file_path:
                chunk_dir = os.path.dirname(chunk_path)
                if tmp_dir is None:
                    tmp_dir = chunk_dir
                try:
                    os.remove(chunk_path)
                except OSError:
                    pass

    # Cleanup temp directory
    if tmp_dir and tmp_dir != os.path.dirname(file_path):
        try:
            os.rmdir(tmp_dir)
        except OSError:
            pass

    if not transcriptions:
        raise RuntimeError(
            f"All {len(chunks)} transcription chunks failed. No audio content could be extracted."
        )

    if failed_chunks:
        logger.warning(
            "[TRANSCRIBER] Partial transcription: %d/%d chunks succeeded.",
            len(transcriptions),
            len(chunks),
        )

    full_text = " ".join(transcriptions)
    logger.info(
        "[TRANSCRIBER] Chunked transcription: %d/%d chunks -> %d chars",
        len(transcriptions), len(chunks), len(full_text),
    )
    return full_text


async def describe_image(image_bytes: bytes, mimetype: str = "image/jpeg") -> str:
    """
    Describe an image using OpenRouter vision model.

    Args:
        image_bytes: Raw image bytes.
        mimetype: MIME type of the image (default: image/jpeg).

    Returns:
        Description text in Portuguese.
    """
    b64_image = base64.b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{mimetype};base64,{b64_image}"

    url = f"{settings.OPENROUTER_BASE_URL}/chat/completions"
    headers = {
        "Authorization": f"Bearer {settings.OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": settings.MODEL_VISION,
        "messages": [
            {
                "role": "system",
                "content": (
                    "Descreva detalhadamente esta imagem em portugues. "
                    "Se houver texto, transcreva."
                ),
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": data_url},
                    }
                ],
            },
        ],
        "max_tokens": 2000,
        "temperature": 0.2,
    }

    async with httpx.AsyncClient(timeout=120.0) as client:
        response = await client.post(url, headers=headers, json=payload)
        response.raise_for_status()

    data = response.json()
    description = data["choices"][0]["message"]["content"].strip()
    logger.info("[TRANSCRIBER] Described image -> %d chars", len(description))
    return description


def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    doc = fitz.open(file_path)
    text_parts = [page.get_text() for page in doc]
    num_pages = len(doc)
    doc.close()
    text = "\n".join(text_parts).strip()
    logger.info("[TRANSCRIBER] PDF text: %d chars from %d pages", len(text), num_pages)
    return text


async def _ocr_pdf_via_vision(file_path: str) -> str:
    """Render PDF pages as images and OCR via vision model."""
    doc = fitz.open(file_path)
    descriptions = []
    max_pages = min(len(doc), 10)

    for i in range(max_pages):
        page = doc[i]
        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes("png")
        desc = await describe_image(img_bytes, "image/png")
        if desc:
            descriptions.append(f"[Pagina {i + 1}]\n{desc}")

    doc.close()
    result = "\n\n".join(descriptions)
    logger.info("[TRANSCRIBER] PDF OCR via vision: %d pages -> %d chars", max_pages, len(result))
    return result


def _extract_docx(file_path: str) -> str:
    """Extract text from DOCX (Word)."""
    from docx import Document

    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    logger.info("[TRANSCRIBER] DOCX text: %d chars", len(text))
    return text


def _extract_xlsx(file_path: str) -> str:
    """Extract text from XLSX (Excel)."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"[Planilha: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    text = "\n".join(parts).strip()
    logger.info("[TRANSCRIBER] XLSX text: %d chars", len(text))
    return text


def _extract_pptx(file_path: str) -> str:
    """Extract text from PPTX (PowerPoint)."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(" | ".join(cells))
        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
    text = "\n\n".join(parts).strip()
    logger.info("[TRANSCRIBER] PPTX text: %d chars from %d slides", len(text), len(prs.slides))
    return text


def _extract_plain_text(file_path: str) -> str:
    """Read file as plain text (TXT, CSV, JSON, XML, HTML, MD, etc.)."""
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                text = f.read(500_000)  # Max 500KB of text
            logger.info("[TRANSCRIBER] Plain text: %d chars (%s)", len(text), enc)
            return text.strip()
        except (UnicodeDecodeError, ValueError):
            continue
    return ""


# Extensions that can be read as plain text
PLAIN_TEXT_EXTENSIONS = {
    ".txt", ".csv", ".json", ".xml", ".html", ".htm", ".md",
    ".yaml", ".yml", ".log", ".ini", ".cfg", ".conf", ".toml",
    ".py", ".js", ".ts", ".sql", ".sh", ".bat", ".css",
}

# MIME types mapped to format
MIME_FORMAT_MAP = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "doc",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xls",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "ppt",
    "text/plain": "txt",
    "text/csv": "csv",
    "application/json": "json",
    "text/html": "html",
    "text/xml": "xml",
    "application/xml": "xml",
}


def _detect_format(file_path: str, mimetype: str | None) -> str:
    """Detect document format from extension and mimetype."""
    ext = os.path.splitext(file_path)[1].lower()

    # Extension-based detection
    ext_map = {
        ".pdf": "pdf",
        ".docx": "docx", ".doc": "doc",
        ".xlsx": "xlsx", ".xls": "xls",
        ".pptx": "pptx", ".ppt": "ppt",
    }
    if ext in ext_map:
        return ext_map[ext]
    if ext in PLAIN_TEXT_EXTENSIONS:
        return "text"

    # MIME-based fallback
    if mimetype:
        mime_lower = mimetype.lower()
        if mime_lower in MIME_FORMAT_MAP:
            return MIME_FORMAT_MAP[mime_lower]
        if mime_lower.startswith("text/"):
            return "text"

    return "unknown"


async def extract_document_text(file_path: str, mimetype: str | None = None) -> str:
    """
    Extract text from any document file.

    Supported formats:
    - PDF: PyMuPDF text extraction, vision OCR fallback for scanned docs
    - DOCX: python-docx (paragraphs + tables)
    - XLSX: openpyxl (all sheets, all rows)
    - PPTX: python-pptx (all slides, shapes + tables)
    - TXT/CSV/JSON/XML/HTML/MD/code: plain text read
    - DOC/XLS/PPT (legacy Office): vision OCR via rendered pages
    - Unknown: skip

    Args:
        file_path: Path to the document file on disk.
        mimetype: MIME type of the document.

    Returns:
        Extracted text content.
    """
    fmt = _detect_format(file_path, mimetype)
    logger.info("[TRANSCRIBER] Document format detected: %s (mime=%s)", fmt, mimetype)

    try:
        if fmt == "pdf":
            text = _extract_pdf(file_path)
            # If scanned PDF (little text), try vision OCR
            if len(text) < 100 and settings.OPENROUTER_API_KEY:
                logger.info("[TRANSCRIBER] PDF has little text (%d chars), trying vision OCR...", len(text))
                ocr_text = await _ocr_pdf_via_vision(file_path)
                return ocr_text if ocr_text else text
            return text

        if fmt == "docx":
            return _extract_docx(file_path)

        if fmt == "xlsx":
            return _extract_xlsx(file_path)

        if fmt == "pptx":
            return _extract_pptx(file_path)

        if fmt == "text":
            return _extract_plain_text(file_path)

        # Legacy Office formats (doc/xls/ppt) - no native parser, try vision
        if fmt in ("doc", "xls", "ppt") and settings.OPENROUTER_API_KEY:
            logger.info("[TRANSCRIBER] Legacy Office format %s, trying PyMuPDF render...", fmt)
            try:
                # PyMuPDF can open some legacy formats
                text = _extract_pdf(file_path)
                if len(text) > 50:
                    return text
                return await _ocr_pdf_via_vision(file_path)
            except Exception:
                logger.warning("[TRANSCRIBER] PyMuPDF can't handle %s format", fmt)
                return ""

        logger.info("[TRANSCRIBER] Unsupported document format: %s", fmt)
        return ""

    except Exception as exc:
        logger.error("[TRANSCRIBER] Document extraction failed for %s: %s", fmt, exc)
        return ""
